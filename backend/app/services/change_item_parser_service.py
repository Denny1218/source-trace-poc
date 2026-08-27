"""Structured change-item parsing from program change-history PPTX slides."""

from __future__ import annotations

import json
import re
from dataclasses import asdict, dataclass, field

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.table import Table

from app.core.logging import get_logger

logger = get_logger()

PARSER_VERSION = 1

META_HEADER_ALIASES: dict[str, tuple[str, ...]] = {
    "item_no": ("번호",),
    "change_title": ("변경사항 정의",),
    "csr_no": ("관련 csr 번호", "관련 csr번호", "csr 번호"),
    "business_background": ("business 관점의 의미/배경", "business 관점의 의미 배경"),
}

DETAIL_HEADER_ALIASES: dict[str, tuple[str, ...]] = {
    "as_is_label": ("as-is", "as is"),
    "to_be_label": ("to-be", "to be"),
    "source_functions": ("소스/함수명", "소스 함수명"),
    "test_cases": ("scenario or case #", "scenario or case", "test case"),
}

NEGATIVE_TABLE_MARKERS = (
    "배포기관 및 f/w",
    "배포기관 및 fw",
    "배포계획",
    "일시 기관",
    "v1 연계항목id",
    "변경 항목 리스트",
    "drb",
    "drb/설계검토",
)

SECTION_DETAIL_LABEL = "변경 설계 상세 내역"

TAG_CURRENT_STATUS = "[현황]"
TAG_AS_IS = "[기존 로직]"
TAG_TO_BE = "[변경 로직]"

SCOPE_PATTERN = re.compile(r"\[([^\]]+)\]")

SOURCE_SPLIT_PATTERN = re.compile(r"\s+-\s+")
# Entry separators: newline / semicolon only. Commas are function-level
# separators inside a "path - func1, func2" entry, so they must not split entries.
MULTI_SOURCE_SPLIT = re.compile(r"[\n;]+")

FILE_EXT_PATTERN = re.compile(
    r"([^\s,;]+\.(?:c|h|cpp|cc|hpp))(?:\s+-\s+|\s+|$)",
    re.IGNORECASE,
)


@dataclass
class SourceFunctionEntry:
    file_path: str | None
    functions: list[str] = field(default_factory=list)
    raw_text: str = ""

    def to_dict(self) -> dict:
        return {
            "file_path": self.file_path,
            "functions": self.functions,
            "raw_text": self.raw_text,
        }


@dataclass
class ChangeItem:
    slide_no: int
    item_no: str | None = None
    change_title: str | None = None
    csr_no: str | None = None
    business_background: str | None = None
    current_status: str | None = None
    as_is: str | None = None
    to_be: str | None = None
    source_functions: list[SourceFunctionEntry] = field(default_factory=list)
    test_cases: list[str] = field(default_factory=list)
    applicable_scopes: list[str] = field(default_factory=list)
    raw_text: str = ""
    template_profile: str | None = None

    def to_dict(self) -> dict:
        data = asdict(self)
        data["source_functions"] = [entry.to_dict() for entry in self.source_functions]
        return data


def normalize_label(text: str | None) -> str:
    if not text:
        return ""
    normalized = re.sub(r"\s+", " ", text.replace("\r", " ").replace("\n", " ").strip())
    return normalized.lower()


def _cell_text(table: Table, row: int, col: int) -> str:
    if row >= len(table.rows) or col >= len(table.columns):
        return ""
    return (table.cell(row, col).text or "").strip()


def _row_labels(table: Table, row: int) -> list[str]:
    return [normalize_label(_cell_text(table, row, col)) for col in range(len(table.columns))]


def _map_row_labels(
    table: Table,
    row: int,
    aliases: dict[str, tuple[str, ...]],
) -> dict[str, int]:
    mapping: dict[str, int] = {}
    for col in range(len(table.columns)):
        label = normalize_label(_cell_text(table, row, col))
        if not label:
            continue
        for field_name, options in aliases.items():
            if field_name in mapping:
                continue
            if any(label == opt or opt in label for opt in options):
                mapping[field_name] = col
    return mapping


def _infer_template_profile(rows: int, cols: int, detail_map: dict[str, int]) -> str:
    if "source_functions" in detail_map:
        source_col = detail_map["source_functions"]
        if rows == 5 and cols == 6 and source_col == 3:
            return "geometry-5x6"
        if rows == 5 and cols == 7 and source_col == 4:
            return "geometry-5x7"
        if rows == 5 and cols == 8 and source_col == 3:
            return "geometry-5x8"
    return f"geometry-{rows}x{cols}"


def _fallback_detail_map(table: Table) -> dict[str, int]:
    rows = len(table.rows)
    cols = len(table.columns)
    detail_row = 3 if rows > 3 else rows - 1
    mapped = _map_row_labels(table, detail_row, DETAIL_HEADER_ALIASES)
    if "source_functions" in mapped:
        return mapped
    # Known geometry fallback (diagnostic only)
    if rows == 5 and cols == 7:
        return {
            "as_is_label": 0,
            "to_be_label": 2,
            "source_functions": 4,
            "test_cases": 6,
        }
    if rows >= 5 and cols >= 6:
        return {
            "as_is_label": 0,
            "to_be_label": 2,
            "source_functions": 3,
            "test_cases": 5,
        }
    return mapped


def is_negative_table(table: Table) -> bool:
    if not table.rows or not table.columns:
        return True
    first = normalize_label(_cell_text(table, 0, 0))
    if any(marker in first for marker in NEGATIVE_TABLE_MARKERS):
        return True
    row0 = " ".join(_row_labels(table, 0))
    if any(marker in row0 for marker in NEGATIVE_TABLE_MARKERS):
        return True
    return False


def _table_has_section_detail(table: Table) -> bool:
    for row in range(len(table.rows)):
        for col in range(len(table.columns)):
            text = normalize_label(_cell_text(table, row, col))
            if SECTION_DETAIL_LABEL.lower() in text:
                return True
    return False


def _table_has_detail_headers(table: Table) -> bool:
    for row in range(len(table.rows)):
        labels = _row_labels(table, row)
        joined = " ".join(labels)
        if "as-is" in joined or "as is" in joined:
            if "to-be" in joined or "to be" in joined or "소스/함수명" in joined:
                return True
    return False


def is_change_item_table(table: Table) -> bool:
    if is_negative_table(table):
        return False
    meta = _map_row_labels(table, 0, META_HEADER_ALIASES)
    if "change_title" not in meta or "item_no" not in meta:
        return False
    if not _table_has_section_detail(table):
        return False
    if not _table_has_detail_headers(table):
        return False
    return True


def find_main_change_item_table(slide: Slide) -> Table | None:
    candidates: list[Table] = []
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            table = shape.table
            if is_change_item_table(table):
                candidates.append(table)
    if not candidates:
        return None
    # Prefer largest matching table
    return max(candidates, key=lambda t: len(t.rows) * len(t.columns))


def is_change_item_slide(slide: Slide) -> bool:
    if not slide.shapes:
        return False
    has_picture_only = all(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes
    )
    if has_picture_only:
        return False
    return find_main_change_item_table(slide) is not None


def _split_tagged_text(text: str) -> tuple[str | None, str | None, str | None]:
    if not text or not text.strip():
        return None, None, None
    raw = text.strip()
    current_status: str | None = None
    as_is: str | None = None
    to_be: str | None = None

    def extract(tag: str, source: str) -> tuple[str | None, str]:
        idx = source.find(tag)
        if idx < 0:
            return None, source
        start = idx + len(tag)
        remainder = source[start:].strip()
        next_tags = [TAG_AS_IS, TAG_TO_BE, TAG_CURRENT_STATUS]
        end = len(remainder)
        for other in next_tags:
            if other == tag:
                continue
            pos = remainder.find(other)
            if pos >= 0:
                end = min(end, pos)
        value = remainder[:end].strip()
        return value or None, source

    if TAG_CURRENT_STATUS in raw:
        current_status, _ = extract(TAG_CURRENT_STATUS, raw)
    if TAG_AS_IS in raw:
        as_is, _ = extract(TAG_AS_IS, raw)
    elif current_status is None and TAG_TO_BE not in raw:
        as_is = raw
    if TAG_TO_BE in raw:
        to_be, _ = extract(TAG_TO_BE, raw)

    if current_status is None and as_is is None and to_be is None:
        return None, None, raw
    return current_status, as_is, to_be


def parse_source_functions(text: str | None) -> list[SourceFunctionEntry]:
    if not text or not text.strip():
        return []
    raw = text.strip()
    entries: list[SourceFunctionEntry] = []

    chunks = [part.strip() for part in MULTI_SOURCE_SPLIT.split(raw) if part.strip()]
    if not chunks:
        chunks = [raw]

    for chunk in chunks:
        if " - " in chunk:
            path_part, func_part = SOURCE_SPLIT_PATTERN.split(chunk, maxsplit=1)
            functions = [f.strip() for f in re.split(r"[,/]", func_part) if f.strip()]
            entries.append(
                SourceFunctionEntry(
                    file_path=path_part.strip() or None,
                    functions=functions,
                    raw_text=chunk,
                )
            )
            continue

        ext_match = FILE_EXT_PATTERN.search(chunk)
        if ext_match:
            file_path = ext_match.group(1).strip()
            remainder = chunk[ext_match.end() :].strip(" -")
            functions = [remainder] if remainder else []
            entries.append(
                SourceFunctionEntry(
                    file_path=file_path,
                    functions=functions,
                    raw_text=chunk,
                )
            )
            continue

        entries.append(SourceFunctionEntry(file_path=None, functions=[], raw_text=chunk))

    return entries


def _extract_applicable_scopes(slide: Slide) -> list[str]:
    scopes: list[str] = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text or ""
        for match in SCOPE_PATTERN.findall(text):
            token = match.strip()
            if token and token not in scopes:
                scopes.append(token)
    return scopes


def _find_detail_content_row(table: Table) -> int:
    for row in range(len(table.rows)):
        labels = _row_labels(table, row)
        if any("as-is" in label or "as is" in label for label in labels):
            return min(row + 1, len(table.rows) - 1)
    return len(table.rows) - 1


def _find_detail_header_row(table: Table) -> int:
    for row in range(len(table.rows)):
        labels = _row_labels(table, row)
        if any("as-is" in label or "as is" in label for label in labels):
            return row
    return 3 if len(table.rows) > 3 else 0


def _value_at(table: Table, row: int, col: int | None) -> str | None:
    if col is None:
        return None
    text = _cell_text(table, row, col)
    return text if text else None


def parse_change_item_from_table(
    table: Table,
    slide_no: int,
    applicable_scopes: list[str] | None = None,
) -> ChangeItem | None:
    if not is_change_item_table(table):
        return None

    meta_map = _map_row_labels(table, 0, META_HEADER_ALIASES)
    detail_header_row = _find_detail_header_row(table)
    detail_map = _map_row_labels(table, detail_header_row, DETAIL_HEADER_ALIASES)
    if "source_functions" not in detail_map:
        detail_map = {**_fallback_detail_map(table), **detail_map}

    content_row = _find_detail_content_row(table)
    profile = _infer_template_profile(len(table.rows), len(table.columns), detail_map)

    item_no = _value_at(table, 1, meta_map.get("item_no"))
    change_title = _value_at(table, 1, meta_map.get("change_title"))
    csr_no = _value_at(table, 1, meta_map.get("csr_no"))
    business_background = _value_at(table, 1, meta_map.get("business_background"))

    as_is_cell = _value_at(table, content_row, detail_map.get("as_is_label"))
    to_be_cell = _value_at(table, content_row, detail_map.get("to_be_label"))
    source_cell = _value_at(table, content_row, detail_map.get("source_functions"))
    test_cell = _value_at(table, content_row, detail_map.get("test_cases"))

    current_status, as_is, to_be_from_as = _split_tagged_text(as_is_cell or "")
    _, _, to_be_from_to = _split_tagged_text(to_be_cell or "")
    to_be = to_be_from_to or to_be_from_as

    source_functions = parse_source_functions(source_cell)
    test_cases: list[str] = []
    if test_cell:
        test_cases = [part.strip() for part in MULTI_SOURCE_SPLIT.split(test_cell) if part.strip()]

    raw_parts = [
        p
        for p in [
            item_no,
            change_title,
            csr_no,
            business_background,
            as_is_cell,
            to_be_cell,
            source_cell,
            test_cell,
        ]
        if p
    ]
    raw_text = "\n".join(raw_parts)

    return ChangeItem(
        slide_no=slide_no,
        item_no=item_no,
        change_title=change_title,
        csr_no=csr_no,
        business_background=business_background,
        current_status=current_status,
        as_is=as_is,
        to_be=to_be,
        source_functions=source_functions,
        test_cases=test_cases,
        applicable_scopes=applicable_scopes or [],
        raw_text=raw_text,
        template_profile=profile,
    )


def parse_change_item_from_slide(slide: Slide, slide_no: int) -> ChangeItem | None:
    table = find_main_change_item_table(slide)
    if table is None:
        return None
    scopes = _extract_applicable_scopes(slide)
    return parse_change_item_from_table(table, slide_no, scopes)


def parse_change_items_from_presentation(presentation: Presentation) -> list[ChangeItem]:
    items: list[ChangeItem] = []
    for index, slide in enumerate(presentation.slides, start=1):
        try:
            item = parse_change_item_from_slide(slide, index)
        except Exception as exc:
            logger.warning(
                "Change item parse failed slide=%s exception_type=%s",
                index,
                type(exc).__name__,
            )
            continue
        if item is not None:
            items.append(item)
    return items


def parse_change_items_from_file(file_path: str) -> list[ChangeItem]:
    presentation = Presentation(file_path)
    return parse_change_items_from_presentation(presentation)


def change_items_to_json(items: list[ChangeItem]) -> str:
    return json.dumps([item.to_dict() for item in items], ensure_ascii=False)


def extract_shape_table_for_test(shape: BaseShape) -> Table | None:
    if getattr(shape, "has_table", False):
        return shape.table
    return None
