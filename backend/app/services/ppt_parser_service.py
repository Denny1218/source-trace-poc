"""PPTX parsing via python-pptx (no DB dependency)."""

from __future__ import annotations

from dataclasses import dataclass, field

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.slide import Slide

from app.core.logging import get_logger
from app.core.ppt_parse_config import FALLBACK_TITLE_MAX_LENGTH

logger = get_logger()


@dataclass
class ParsedSlide:
    slide_number: int
    title: str | None
    content: str


@dataclass
class ParsedPresentation:
    slide_count: int
    slides: list[ParsedSlide] = field(default_factory=list)


class PptParseError(Exception):
    def __init__(self, message: str):
        self.message = message
        super().__init__(message)


def _normalize_paragraph_text(text: str) -> str:
    return text.strip()


def _normalize_content(text: str) -> str:
    lines = [line.rstrip() for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")]
    normalized: list[str] = []
    blank_run = 0
    for line in lines:
        if not line.strip():
            blank_run += 1
            if blank_run <= 2:
                normalized.append("")
            continue
        blank_run = 0
        normalized.append(line)
    while normalized and normalized[-1] == "":
        normalized.pop()
    return "\n".join(normalized)


def _extract_text_frame(shape: BaseShape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    paragraphs: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = _normalize_paragraph_text(paragraph.text)
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def _extract_table_text(shape: BaseShape) -> str:
    if not getattr(shape, "has_table", False):
        return ""
    rows: list[str] = []
    for row in shape.table.rows:
        cells: list[str] = []
        for cell in row.cells:
            cell_text = _normalize_paragraph_text(cell.text)
            cells.append(cell_text)
        rows.append(" | ".join(cells))
    return "\n".join(rows)


def _is_group_shape(shape: BaseShape) -> bool:
    return isinstance(shape, GroupShape) or shape.shape_type == MSO_SHAPE_TYPE.GROUP


def _extract_shape_texts(shape: BaseShape, slide_number: int) -> list[str]:
    """Extract text from a shape using group → table → text frame order."""
    if _is_group_shape(shape):
        texts: list[str] = []
        for child in shape.shapes:
            texts.extend(_extract_shape_texts(child, slide_number))
        return texts

    if getattr(shape, "has_table", False):
        table_text = _extract_table_text(shape)
        return [table_text] if table_text else []

    if getattr(shape, "has_text_frame", False):
        frame_text = _extract_text_frame(shape)
        return [frame_text] if frame_text else []

    # Unsupported shape types (SmartArt, Chart, Picture, etc.) — skip without failing
    logger.debug(
        "PPT unsupported shape skipped slide=%s shape_type=%s",
        slide_number,
        shape.shape_type,
    )
    return []


def _meaningful_lines(text: str) -> list[str]:
    return [line.strip() for line in text.split("\n") if line.strip()]


def resolve_slide_title(slide: Slide, content: str) -> str | None:
    """Title policy: Title Placeholder → first meaningful line → None."""
    title_shape = slide.shapes.title
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        title_text = _normalize_paragraph_text(title_shape.text)
        if title_text:
            return title_text[:FALLBACK_TITLE_MAX_LENGTH]

    for shape in slide.shapes:
        if shape == title_shape:
            continue
        for block in _extract_shape_texts(shape, 0):
            for line in _meaningful_lines(block):
                return line[:FALLBACK_TITLE_MAX_LENGTH]

    for line in _meaningful_lines(content):
        return line[:FALLBACK_TITLE_MAX_LENGTH]

    return None


def _parse_slide(slide: Slide, slide_number: int) -> ParsedSlide:
    content_blocks: list[str] = []
    for shape in slide.shapes:
        content_blocks.extend(_extract_shape_texts(shape, slide_number))

    content = _normalize_content("\n\n".join(block for block in content_blocks if block))
    title = resolve_slide_title(slide, content)
    return ParsedSlide(slide_number=slide_number, title=title, content=content)


def parse_pptx_file(file_path: str) -> ParsedPresentation:
    """Parse a PPTX file and return slide text (1-based slide numbers)."""
    try:
        presentation = Presentation(file_path)
    except Exception as exc:
        raise PptParseError(f"PPTX 열기 실패: {file_path}") from exc

    slides: list[ParsedSlide] = []
    for index, slide in enumerate(presentation.slides, start=1):
        try:
            slides.append(_parse_slide(slide, index))
        except Exception as exc:
            logger.warning(
                "PPT slide parse failed file=%s slide=%s exception_type=%s",
                file_path,
                index,
                type(exc).__name__,
            )
            slides.append(ParsedSlide(slide_number=index, title=None, content=""))

    return ParsedPresentation(slide_count=len(slides), slides=slides)


def extract_shape_texts_for_test(shape: BaseShape, slide_number: int = 1) -> list[str]:
    """Expose shape extraction for unit tests."""
    return _extract_shape_texts(shape, slide_number)
